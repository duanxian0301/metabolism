from pathlib import Path
from xml.sax.saxutils import escape

from PIL import Image, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(r"D:\codex\GenomicSEM\metabolic")
FIG = ROOT / "figures"
SOURCE = next(FIG.glob("ChatGPT Image*.png"))
OUT = FIG / "figure1_editable_outputs"
OUT.mkdir(exist_ok=True)

W_IN, H_IN = 10.55, 14.91

COLORS = {
    "blue": "154D9B",
    "blue_fill": "EEF5FF",
    "green": "2E7D32",
    "green_fill": "F2FBF1",
    "purple": "6B3FA0",
    "purple_fill": "F7F2FF",
    "orange": "E85D14",
    "orange_fill": "FFF6EF",
    "gold": "C98500",
    "gold_fill": "FFF9EA",
    "red": "E24A2A",
    "gray": "4A5568",
    "light": "FAFCFF",
    "line": "5F6F7D",
    "black": "111111",
}


def rgb(hexstr):
    hexstr = hexstr.replace("#", "")
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


def add_text(slide, text, x, y, w, h, size=12, bold=False, color="black", align="center"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(COLORS[color])
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    return box


def add_box(slide, x, y, w, h, stroke, fill, radius=True, lw=1.2):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(COLORS[fill])
    s.line.color.rgb = rgb(COLORS[stroke])
    s.line.width = Pt(lw)
    return s


def add_arrow(slide, x1, y1, x2, y2, color="line", width=1.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(COLORS[color])
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_step_badge(slide, n, x, y, color):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.38), Inches(0.38))
    c.fill.solid()
    c.fill.fore_color.rgb = rgb(COLORS[color])
    c.line.color.rgb = rgb(COLORS[color])
    add_text(slide, str(n), x, y + 0.005, 0.38, 0.35, size=16, bold=True, color="light")


def add_icon_label(slide, icon, label, x, y, w, color):
    add_text(slide, icon, x, y, w, 0.34, size=23, bold=False, color=color)
    add_text(slide, label, x, y + 0.38, w, 0.56, size=10.2, bold=True, color="black")


def make_pptx():
    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(COLORS["light"])

    add_text(slide, "Figure 1. Overview of study design and analytical workflow", 1.95, 0.08, 6.65, 0.36, 18, True)

    # Section 1
    add_box(slide, 0.12, 0.55, 10.31, 1.78, "blue", "light", lw=1.3)
    add_step_badge(slide, 1, 0.22, 0.68, "blue")
    add_text(slide, "Data sources", 0.78, 0.65, 1.6, 0.32, 11.5, True, "blue", "left")
    add_text(slide, "▰\n▰\n▰", 0.22, 1.20, 0.56, 0.70, 18, True, "blue")
    add_box(slide, 1.02, 1.08, 2.95, 1.14, "blue", "light")
    add_text(slide, "NMR metabolite GWAS (meta_EUR)\n249 circulating metabolic traits\nEuropean ancestry\nn = 599,249", 1.23, 1.16, 2.50, 0.88, 9.2, True)
    add_box(slide, 4.08, 1.08, 2.42, 1.14, "blue", "light")
    add_text(slide, "External metabolite GWAS\n233 NMR biomarkers\nup to n = 136,016", 4.35, 1.25, 1.88, 0.64, 9.2, True)
    add_box(slide, 6.60, 1.08, 3.72, 1.14, "blue", "light")
    add_text(slide, "Neurodegenerative disease GWAS\nAD: 39,106 cases + 46,828 proxy cases + 401,577 controls\nPD: 63,555 cases + 17,700 proxy cases + 1,746,386 controls\nLBD: 2,591 cases + 4,027 controls", 7.20, 1.14, 2.90, 0.92, 7.9, True)
    add_text(slide, "DNA", 1.10, 1.28, 0.42, 0.40, 11, True, "blue")
    add_text(slide, "●─●", 4.20, 1.35, 0.48, 0.25, 14, True, "blue")
    add_text(slide, "🧠", 6.68, 1.35, 0.45, 0.35, 18, False, "blue")

    add_arrow(slide, 5.28, 2.33, 5.28, 2.70, "black", 1.8)

    # Section 2
    add_box(slide, 0.12, 2.68, 10.31, 2.90, "green", "green_fill", lw=1.2)
    add_step_badge(slide, 2, 0.22, 2.82, "green")
    add_text(slide, "Trait harmonization\nand factor\nconstruction", 0.18, 3.82, 1.12, 1.00, 12, True, "green")
    stage2 = [
        ("☑", "Summary-statistic\nQC and\nharmonization"),
        ("▽", "Retain non-\nproportion traits\nwith h2 Z > 4"),
        ("▂▅▇", "112\ntraits"),
        ("◔", "Split into lipid\nand non-lipid\ndomains"),
        ("⌁", "Staged reduction\nto final 8 + 8\nindicators"),
        ("⌬", "Genomic SEM\nfactor\nmodeling"),
    ]
    x = 1.40
    widths = [1.50, 1.35, 0.95, 1.30, 1.45, 1.35]
    for i, ((ic, lab), ww) in enumerate(zip(stage2, widths)):
        add_box(slide, x, 2.86, ww, 0.92, "green", "light")
        add_icon_label(slide, ic, lab, x + 0.08, 2.98, ww - 0.16, "green")
        if i < len(stage2) - 1:
            add_arrow(slide, x + ww + 0.03, 3.32, x + ww + 0.18, 3.32, "black", 1.0)
        x += ww + 0.25
    add_box(slide, 1.60, 4.22, 3.95, 1.20, "green", "light")
    add_text(slide, "Lipid factors", 3.04, 4.30, 1.08, 0.22, 10.5, True, "green")
    for x, lab in [(1.72, "TG-rich /\nVLDL-remodeling\naxis"), (3.10, "HDL-core\naxis"), (4.20, "Cholesteryl-ester /\nstructural-lipid\naxis")]:
        add_box(slide, x, 4.62, 1.28, 0.68, "green", "green_fill")
        add_text(slide, lab, x + 0.04, 4.68, 1.20, 0.48, 8.7, True, "green")
    add_box(slide, 5.86, 4.22, 3.95, 1.20, "green", "light")
    add_text(slide, "Non-lipid factors", 7.16, 4.30, 1.34, 0.22, 10.5, True, "green")
    for x, lab in [(5.98, "Ketone-body\naxis"), (7.18, "Amino-acid\naxis"), (8.42, "Energy-bridge\naxis")]:
        add_box(slide, x, 4.62, 1.12, 0.68, "green", "green_fill")
        add_text(slide, lab, x + 0.04, 4.68, 1.04, 0.48, 8.7, True, "green")

    add_arrow(slide, 5.28, 5.58, 5.28, 5.93, "black", 1.8)

    # Section 3
    add_box(slide, 0.12, 5.92, 10.31, 1.82, "purple", "purple_fill", lw=1.2)
    add_step_badge(slide, 3, 0.22, 6.07, "purple")
    add_text(slide, "Factor GWAS\ngeneration and\nvalidation", 0.20, 6.70, 1.18, 0.70, 11.5, True, "purple")
    items3 = [
        ("DNA+", "Generate six\nlatent factor\nGWAS"),
        ("▂▅▇", "Estimate\neffective\nsample size"),
        ("↗", "Univariate and\nbivariate LDSC\nQC"),
        ("✓", "Internal and\nexternal\nmetabolite\nvalidation"),
        ("↻", "Leave-one-out\nsensitivity\nanalysis"),
        ("∩", "Q_SNP\nheterogeneity\nannotation"),
    ]
    x = 1.58
    for ic, lab in items3:
        add_icon_label(slide, ic, lab, x, 6.08, 1.22, "purple")
        x += 1.36

    add_arrow(slide, 5.28, 7.74, 5.28, 8.10, "black", 1.8)

    # Section 4
    add_box(slide, 0.12, 8.08, 10.31, 1.95, "orange", "orange_fill", lw=1.2)
    add_step_badge(slide, 4, 0.22, 8.22, "orange")
    add_text(slide, "Cross-disease\nscreening", 0.20, 8.95, 1.25, 0.50, 11.5, True, "orange")
    add_text(slide, "6 metabolic factors x 3 diseases (AD, PD, LBD)\nBivariate LDSC screen with FDR correction\n↓\nMiXeR for prioritized branches", 3.35, 8.17, 3.85, 0.86, 10.0, True)
    branches = [
        (1.85, "HDL-core <-> AD", "blue", "blue_fill"),
        (4.05, "Ketone-body <-> PD", "purple", "purple_fill"),
        (6.85, "TG-rich / VLDL-remodeling <-> PD", "orange", "orange_fill"),
    ]
    for x, lab, stroke, fill in branches:
        add_box(slide, x, 9.10, 1.95 if x < 6 else 3.25, 0.45, stroke, fill)
        add_text(slide, lab, x + 0.10, 9.17, (1.75 if x < 6 else 3.05), 0.22, 9.4, True, stroke)
    add_box(slide, 3.35, 9.68, 4.10, 0.25, "orange", "light")
    add_text(slide, "All LBD comparisons were not FDR-significant", 3.42, 9.695, 3.95, 0.16, 8.6, True)

    add_arrow(slide, 5.28, 10.03, 5.28, 10.38, "black", 1.8)

    # Section 5
    add_box(slide, 0.12, 10.38, 10.31, 1.82, "gold", "gold_fill", lw=1.2)
    add_step_badge(slide, 5, 0.22, 10.53, "gold")
    add_text(slide, "Locus, gene,\ncellular, and\ntarget analyses", 0.18, 11.12, 1.20, 0.70, 11, True, "gold")
    pipeline = [
        ("⌕", "conjFDR\nshared-locus\ndiscovery"),
        ("∩", "coloc +\nPWCoCo"),
        ("Gene prioritization", "FUMA   cTWAS   bulk-brain SMR\nGTEx brain      brain cell-type SMR"),
        ("✣", "scPagwas\nprojection in\nAD and PD\nsingle-cell\natlases"),
        ("⟱", "Targeted\nscTenifoldKnk\nvirtual\nknockout"),
        ("Target annotation", "Open Targets   DGIdb   Human Protein Atlas\nchemical probes   protein interactions   AlphaFold"),
    ]
    x = 1.30
    widths = [0.88, 1.02, 2.02, 1.05, 1.05, 2.08]
    for i, (node, ww) in enumerate(zip(pipeline, widths)):
        add_box(slide, x, 10.55, ww, 1.45, "gold", "light")
        if i in [2, 5]:
            add_text(slide, node[0], x + 0.08, 10.66, ww - 0.16, 0.22, 8.8, True)
            add_text(slide, node[1], x + 0.10, 10.98, ww - 0.20, 0.68, 7.4, True)
        else:
            add_icon_label(slide, node[0], node[1], x + 0.05, 10.68, ww - 0.10, "gold")
        if i < len(widths) - 1:
            add_arrow(slide, x + ww + 0.02, 11.26, x + ww + 0.14, 11.26, "black", 1.0)
        x += ww + 0.24

    add_arrow(slide, 5.28, 12.20, 5.28, 12.55, "black", 1.8)

    # Section 6
    add_box(slide, 0.12, 12.55, 10.31, 2.20, "blue", "blue_fill", lw=1.2)
    add_step_badge(slide, 6, 0.22, 12.72, "blue")
    add_text(slide, "Main biological\nand translational\noutputs", 0.20, 13.40, 1.18, 0.68, 11.2, True, "blue")
    out_nodes = [
        (1.55, 12.72, 2.55, "HDL-core / AD\nPericyte-enriched\nVascular / support pathways", "blue"),
        (4.22, 12.72, 3.00, "Ketone-body / PD\nDistributed neuronal-glial enrichment\nEnergy-adaptation programs", "purple"),
        (7.35, 12.72, 2.92, "TG-rich / VLDL-remodeling / PD\nPericyte-enriched\nLysosomal and membrane-\ntrafficking programs", "orange"),
    ]
    for x, y, w, lab, col in out_nodes:
        add_box(slide, x, y, w, 0.95, col, "light")
        add_text(slide, lab, x + 0.12, y + 0.12, w - 0.24, 0.68, 8.8, True, col)
        add_arrow(slide, x + w / 2, y + 0.96, x + w / 2, 14.15, col, 1.4)
    add_box(slide, 1.86, 14.12, 7.78, 0.50, "blue", "light")
    add_text(slide, "Prioritized candidate targets: GRK4, PRMT7, TMEM175", 3.45, 14.18, 4.20, 0.28, 13.5, True)

    pptx_path = OUT / "figure1_workflow_editable.pptx"
    prs.save(pptx_path)
    return pptx_path


def make_high_res():
    im = Image.open(SOURCE).convert("RGB")
    scale = 4
    up = im.resize((im.width * scale, im.height * scale), Image.Resampling.LANCZOS)
    sharp = up.filter(ImageFilter.UnsharpMask(radius=1.2, percent=135, threshold=3))
    png = OUT / "figure1_workflow_4x_600dpi.png"
    tif = OUT / "figure1_workflow_4x_600dpi.tif"
    tiff = OUT / "figure1_workflow_4x_600dpi.tiff"
    sharp_png = OUT / "figure1_workflow_4x_600dpi_sharpened.png"
    sharp_tiff = OUT / "figure1_workflow_4x_600dpi_sharpened.tiff"
    up.save(png, dpi=(600, 600), optimize=True)
    up.save(tif, dpi=(600, 600), compression="tiff_lzw")
    up.save(tiff, dpi=(600, 600), compression="tiff_lzw")
    sharp.save(sharp_png, dpi=(600, 600), optimize=True)
    sharp.save(sharp_tiff, dpi=(600, 600), compression="tiff_lzw")
    return png, tif, tiff, sharp_png, sharp_tiff


def make_svg():
    pptx_path = OUT / "figure1_workflow_vector.svg"
    # Lightweight vector companion: editable text and shapes, same content as the PPTX.
    w, h = 1055, 1491
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}" viewBox="0 0 {w} {h}">',
        "<style>text{font-family:Arial,sans-serif}.title{font-weight:700;font-size:28px}.stage{font-weight:700;font-size:17px}.body{font-weight:700;font-size:14px}.small{font-weight:700;font-size:12px}.num{font-weight:700;font-size:25px;fill:white}</style>",
        f'<rect width="{w}" height="{h}" fill="#{COLORS["light"]}"/>',
        '<text class="title" x="205" y="36">Figure 1. Overview of study design and analytical workflow</text>',
    ]

    def rect(x, y, ww, hh, stroke, fill, sw=1.4):
        parts.append(f'<rect x="{x}" y="{y}" width="{ww}" height="{hh}" rx="12" fill="#{COLORS[fill]}" stroke="#{COLORS[stroke]}" stroke-width="{sw}"/>')

    def text(lines, x, y, klass="small", color="black", anchor="middle", gap=17):
        for i, line in enumerate(lines.split("\n")):
            parts.append(f'<text class="{klass}" x="{x}" y="{y+i*gap}" text-anchor="{anchor}" fill="#{COLORS[color]}">{escape(line)}</text>')

    def badge(n, x, y, color):
        parts.append(f'<circle cx="{x}" cy="{y}" r="20" fill="#{COLORS[color]}"/>')
        text(str(n), x, y + 8, "num", "light")

    sections = [
        (10, 52, 1035, 180, "blue", "light", 1, "Data sources"),
        (10, 266, 1035, 290, "green", "green_fill", 2, "Trait harmonization\nand factor construction"),
        (10, 588, 1035, 182, "purple", "purple_fill", 3, "Factor GWAS\ngeneration and validation"),
        (10, 802, 1035, 195, "orange", "orange_fill", 4, "Cross-disease\nscreening"),
        (10, 1030, 1035, 182, "gold", "gold_fill", 5, "Locus, gene,\ncellular, and\ntarget analyses"),
        (10, 1244, 1035, 235, "blue", "blue_fill", 6, "Main biological\nand translational\noutputs"),
    ]
    for x, y, ww, hh, stroke, fill, n, label in sections:
        rect(x, y, ww, hh, stroke, fill)
        badge(n, x + 30, y + 32, stroke)
        text(label, x + 95, y + 38 if n == 1 else y + 120, "stage", stroke, "start", 21)

    # Data source boxes
    for x, ww, lab in [
        (102, 292, "NMR metabolite GWAS (meta_EUR)\n249 circulating metabolic traits\nEuropean ancestry\nn = 599,249"),
        (406, 240, "External metabolite GWAS\n233 NMR biomarkers\nup to n = 136,016"),
        (656, 375, "Neurodegenerative disease GWAS\nAD: 39,106 cases + 46,828 proxy cases + 401,577 controls\nPD: 63,555 cases + 17,700 proxy cases + 1,746,386 controls\nLBD: 2,591 cases + 4,027 controls"),
    ]:
        rect(x, 106, ww, 114, "blue", "light")
        text(lab, x + ww / 2, 132, "small", "black", "middle", 20)

    # Section summaries
    text("Summary-statistic QC → h2 Z > 4 → 112 traits → lipid/non-lipid split → final 8 + 8 indicators → Genomic SEM", 590, 318, "body", "black")
    rect(162, 418, 392, 120, "green", "light")
    text("Lipid factors\nTG-rich/VLDL-remodeling axis     HDL-core axis     Cholesteryl-ester/structural-lipid axis", 358, 444, "small", "green", "middle", 28)
    rect(584, 418, 388, 120, "green", "light")
    text("Non-lipid factors\nKetone-body axis        Amino-acid axis        Energy-bridge axis", 778, 444, "small", "green", "middle", 28)
    text("Generate six latent factor GWAS     Estimate effective sample size     LDSC QC     metabolite validation     leave-one-out     Q_SNP annotation", 600, 690, "body", "black")
    text("6 metabolic factors x 3 diseases (AD, PD, LBD)\nBivariate LDSC screen with FDR correction\n↓\nMiXeR for prioritized branches", 530, 842, "body", "black", "middle", 22)
    for x, col, lab in [(182, "blue", "HDL-core <-> AD"), (400, "purple", "Ketone-body <-> PD"), (674, "orange", "TG-rich / VLDL-remodeling <-> PD")]:
        rect(x, 908, 200 if x < 600 else 320, 46, col, "light")
        text(lab, x + (100 if x < 600 else 160), 936, "body", col)
    rect(330, 962, 405, 26, "orange", "light")
    text("All LBD comparisons were not FDR-significant", 532, 980, "small", "black")
    text("conjFDR shared-locus discovery  →  coloc + PWCoCo  →  gene prioritization  →  scPagwas projection  →  targeted scTenifoldKnk  →  target annotation", 610, 1132, "small", "black")
    for x, col, lab in [
        (155, "blue", "HDL-core / AD\nPericyte-enriched\nVascular / support pathways"),
        (418, "purple", "Ketone-body / PD\nDistributed neuronal-glial enrichment\nEnergy-adaptation programs"),
        (735, "orange", "TG-rich / VLDL-remodeling / PD\nPericyte-enriched\nLysosomal and membrane-trafficking programs"),
    ]:
        rect(x, 1260, 250 if x != 418 else 300, 94, col, "light")
        text(lab, x + (125 if x != 418 else 150), 1291, "small", col, "middle", 20)
    rect(186, 1418, 778, 50, "blue", "light")
    text("Prioritized candidate targets: GRK4, PRMT7, TMEM175", 575, 1452, "stage", "black")

    # Main down arrows
    for y1, y2 in [(232, 266), (556, 588), (770, 802), (997, 1030), (1212, 1244)]:
        parts.append(f'<line x1="528" y1="{y1}" x2="528" y2="{y2}" stroke="#111" stroke-width="4" marker-end="url(#arrow)"/>')

    parts.insert(2, '<defs><marker id="arrow" markerWidth="10" markerHeight="10" refX="8" refY="3" orient="auto"><path d="M0,0 L8,3 L0,6 Z" fill="#111"/></marker></defs>')
    parts.append("</svg>")
    pptx_path.write_text("\n".join(parts), encoding="utf-8")
    return pptx_path


if __name__ == "__main__":
    pptx = make_pptx()
    png, tif, tiff, sharp_png, sharp_tiff = make_high_res()
    svg = make_svg()
    print("source", SOURCE)
    print("pptx", pptx)
    print("png", png)
    print("tif", tif)
    print("tiff", tiff)
    print("sharp_png", sharp_png)
    print("sharp_tiff", sharp_tiff)
    print("svg", svg)
